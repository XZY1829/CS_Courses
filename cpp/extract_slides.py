"""Extract text from PPTX and PPT files in the cpp folder."""
import os
import sys

def extract_pptx(filepath, output_path):
    """Extract text from .pptx files using python-pptx."""
    from pptx import Presentation
    prs = Presentation(filepath)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n{'='*60}")
        lines.append(f"=== Slide {i} ===")
        lines.append(f"{'='*60}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    lines.append(row_text)
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))
    print(f"  [PPTX] {os.path.basename(filepath)} -> {len(prs.slides)} slides")

def extract_ppt(filepath, output_path):
    """Extract text from .ppt files using win32com (PowerPoint COM)."""
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()
    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        abs_path = os.path.abspath(filepath)
        presentation = ppt_app.Presentations.Open(abs_path, ReadOnly=True, WithWindow=False)
        lines = []
        for i, slide in enumerate(presentation.Slides, 1):
            lines.append(f"\n{'='*60}")
            lines.append(f"=== Slide {i} ===")
            lines.append(f"{'='*60}\n")
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    for para in shape.TextFrame.TextRange.Paragraphs():
                        text = para.Text.strip()
                        if text:
                            lines.append(text)
                if shape.HasTable:
                    table = shape.Table
                    for r in range(1, table.Rows.Count + 1):
                        row_text = []
                        for c in range(1, table.Columns.Count + 1):
                            row_text.append(table.Cell(r, c).Shape.TextFrame.TextRange.Text.strip())
                        lines.append(" | ".join(row_text))
        presentation.Close()
        ppt_app.Quit()
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        print(f"  [PPT]  {os.path.basename(filepath)} -> {len(lines)} lines")
    except Exception as e:
        print(f"  [ERROR] {os.path.basename(filepath)}: {e}")
        try:
            ppt_app.Quit()
        except:
            pass
    finally:
        pythoncom.CoUninitialize()

def main():
    cpp_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(cpp_dir, "_extracted")
    os.makedirs(output_dir, exist_ok=True)
    
    files = sorted([f for f in os.listdir(cpp_dir) 
                    if f.endswith(('.pptx', '.ppt')) and not f.startswith('~')])
    
    print(f"Found {len(files)} presentation files:")
    for f in files:
        print(f"  - {f}")
    print()
    
    for filename in files:
        filepath = os.path.join(cpp_dir, filename)
        base = os.path.splitext(filename)[0]
        output_path = os.path.join(output_dir, f"{base}.txt")
        
        if filename.endswith('.pptx'):
            extract_pptx(filepath, output_path)
        else:
            extract_ppt(filepath, output_path)
    
    print(f"\nAll extracted to: {output_dir}")

if __name__ == '__main__':
    main()
